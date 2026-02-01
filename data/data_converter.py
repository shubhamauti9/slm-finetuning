"""
Document to Training Data Converter
====================================
Converts DOCX, Excel, PDF, PPT, and TXT files into LLM fine-tuning training data.

Usage:
    from data_converter import DocumentConverter
    
    converter = DocumentConverter()
    training_data = converter.convert_file("path/to/file.docx")
    converter.save_json(training_data, "output.json")
"""

import json
import re
import os
from typing import List, Dict, Optional, Tuple
from pathlib import Path


class DocumentConverter:
    """Converts various document formats to training data JSON."""
    
    def __init__(self):
        self.supported_formats = ['.docx', '.xlsx', '.xls', '.pdf', '.pptx', '.ppt', '.txt']
    
    def convert_file(self, file_path: str, strategy: str = "auto") -> List[Dict]:
        """
        Convert a file to training data format.
        
        Args:
            file_path: Path to the file
            strategy: Extraction strategy - 'auto', 'faq', 'heading', 'chunk', 'table'
        
        Returns:
            List of training examples with 'instruction', 'input', 'output' keys
        """
        ext = Path(file_path).suffix.lower()
        
        if ext == '.docx':
            return self._convert_docx(file_path, strategy)
        elif ext in ['.xlsx', '.xls']:
            return self._convert_excel(file_path, strategy)
        elif ext == '.pdf':
            return self._convert_pdf(file_path, strategy)
        elif ext in ['.pptx', '.ppt']:
            return self._convert_pptx(file_path, strategy)
        elif ext == '.txt':
            return self._convert_txt(file_path, strategy)
        else:
            raise ValueError(f"Unsupported format: {ext}")
    
    def convert_files(self, file_paths: List[str], strategy: str = "auto") -> List[Dict]:
        """Convert multiple files and combine results."""
        all_data = []
        for path in file_paths:
            try:
                data = self.convert_file(path, strategy)
                all_data.extend(data)
                print(f"✅ Converted {path}: {len(data)} examples")
            except Exception as e:
                print(f"❌ Failed to convert {path}: {e}")
        return all_data
    
    def save_json(self, data: List[Dict], output_path: str):
        """Save training data to JSON file."""
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(data, f, indent=4, ensure_ascii=False)
        print(f"✅ Saved {len(data)} examples to {output_path}")
    
    # ==================== DOCX Conversion ====================
    
    def _convert_docx(self, file_path: str, strategy: str) -> List[Dict]:
        """Convert DOCX file to training data."""
        try:
            from docx import Document
        except ImportError:
            raise ImportError("Install python-docx: pip install python-docx")
        
        doc = Document(file_path)
        text_content = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                # Check if it's a heading
                if para.style.name.startswith('Heading'):
                    text_content.append(('heading', para.text.strip()))
                else:
                    text_content.append(('text', para.text.strip()))
        
        # Also extract tables
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            if table_data:
                text_content.append(('table', table_data))
        
        return self._apply_strategy(text_content, strategy)
    
    # ==================== Excel Conversion ====================
    
    def _convert_excel(self, file_path: str, strategy: str) -> List[Dict]:
        """Convert Excel file to training data."""
        try:
            import pandas as pd
        except ImportError:
            raise ImportError("Install pandas: pip install pandas openpyxl")
        
        df = pd.read_excel(file_path)
        training_data = []
        
        # Try to detect Q&A columns
        columns = df.columns.tolist()
        question_col = None
        answer_col = None
        
        # Look for question/answer column patterns
        for col in columns:
            col_lower = col.lower()
            if any(q in col_lower for q in ['question', 'query', 'q', 'instruction', 'prompt']):
                question_col = col
            if any(a in col_lower for a in ['answer', 'response', 'a', 'output', 'reply']):
                answer_col = col
        
        if question_col and answer_col:
            # Direct Q&A mapping
            for _, row in df.iterrows():
                q = str(row[question_col]).strip()
                a = str(row[answer_col]).strip()
                if q and a and q != 'nan' and a != 'nan':
                    training_data.append({
                        "instruction": q,
                        "input": "",
                        "output": a
                    })
        else:
            # Use first column as question, second as answer
            if len(columns) >= 2:
                for _, row in df.iterrows():
                    q = str(row[columns[0]]).strip()
                    a = str(row[columns[1]]).strip()
                    if q and a and q != 'nan' and a != 'nan':
                        training_data.append({
                            "instruction": q,
                            "input": "",
                            "output": a
                        })
        
        return training_data
    
    # ==================== PDF Conversion ====================
    
    def _convert_pdf(self, file_path: str, strategy: str) -> List[Dict]:
        """Convert PDF file to training data."""
        try:
            import fitz  # PyMuPDF
        except ImportError:
            raise ImportError("Install PyMuPDF: pip install PyMuPDF")
        
        doc = fitz.open(file_path)
        text_content = []
        
        for page in doc:
            text = page.get_text()
            if text.strip():
                # Simple paragraph splitting
                paragraphs = text.split('\n\n')
                for para in paragraphs:
                    para = para.strip()
                    if para:
                        text_content.append(('text', para))
        
        doc.close()
        return self._apply_strategy(text_content, strategy)
    
    # ==================== PowerPoint Conversion ====================
    
    def _convert_pptx(self, file_path: str, strategy: str) -> List[Dict]:
        """Convert PowerPoint file to training data."""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("Install python-pptx: pip install python-pptx")
        
        prs = Presentation(file_path)
        training_data = []
        
        for slide in prs.slides:
            title = ""
            content = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if not text:
                        continue
                    
                    # First text box is usually the title
                    if shape.has_text_frame:
                        if shape == slide.shapes[0] or (hasattr(shape, 'is_placeholder') and shape.is_placeholder):
                            if not title:
                                title = text
                                continue
                    
                    content.append(text)
            
            if title and content:
                training_data.append({
                    "instruction": title,
                    "input": "",
                    "output": "\n".join(content)
                })
        
        return training_data
    
    # ==================== TXT Conversion ====================
    
    def _convert_txt(self, file_path: str, strategy: str) -> List[Dict]:
        """Convert TXT file to training data."""
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        text_content = []
        
        # Try to detect structure
        lines = content.split('\n')
        current_heading = None
        current_text = []
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # Check for heading patterns
            if self._is_heading(line):
                if current_heading and current_text:
                    text_content.append(('heading', current_heading))
                    text_content.append(('text', '\n'.join(current_text)))
                current_heading = self._clean_heading(line)
                current_text = []
            else:
                current_text.append(line)
        
        # Don't forget last section
        if current_heading and current_text:
            text_content.append(('heading', current_heading))
            text_content.append(('text', '\n'.join(current_text)))
        elif current_text:
            # No headings found, treat as plain text
            text_content.append(('text', '\n'.join(current_text)))
        
        return self._apply_strategy(text_content, strategy)
    
    # ==================== Extraction Strategies ====================
    
    def _apply_strategy(self, content: List[Tuple], strategy: str) -> List[Dict]:
        """Apply extraction strategy to parsed content."""
        if strategy == "auto":
            strategy = self._detect_best_strategy(content)
        
        if strategy == "faq":
            return self._extract_faq(content)
        elif strategy == "heading":
            return self._extract_heading_content(content)
        elif strategy == "chunk":
            return self._extract_chunks(content)
        else:
            # Default: try FAQ first, then heading, then chunk
            result = self._extract_faq(content)
            if not result:
                result = self._extract_heading_content(content)
            if not result:
                result = self._extract_chunks(content)
            return result
    
    def _detect_best_strategy(self, content: List[Tuple]) -> str:
        """Auto-detect the best extraction strategy."""
        text = ' '.join([c[1] if isinstance(c[1], str) else str(c[1]) for c in content])
        
        # Check for FAQ patterns
        faq_patterns = [
            r'\bQ:\s*', r'\bA:\s*',
            r'\bQuestion:\s*', r'\bAnswer:\s*',
            r'\bQ\.\s*', r'\bA\.\s*',
            r'\?\s*\n',  # Question followed by newline
        ]
        faq_score = sum(len(re.findall(p, text, re.IGNORECASE)) for p in faq_patterns)
        
        # Check for headings
        heading_count = sum(1 for c in content if c[0] == 'heading')
        
        if faq_score > 5:
            return "faq"
        elif heading_count > 3:
            return "heading"
        else:
            return "chunk"
    
    def _extract_faq(self, content: List[Tuple]) -> List[Dict]:
        """Extract FAQ-style Q&A pairs."""
        training_data = []
        text = '\n'.join([c[1] if isinstance(c[1], str) else str(c[1]) for c in content])
        
        # Pattern 1: Q: ... A: ...
        pattern1 = r'Q[:\.]?\s*(.+?)\s*A[:\.]?\s*(.+?)(?=Q[:\.]?\s*|$)'
        matches = re.findall(pattern1, text, re.DOTALL | re.IGNORECASE)
        for q, a in matches:
            q, a = q.strip(), a.strip()
            if q and a:
                training_data.append({
                    "instruction": q,
                    "input": "",
                    "output": a
                })
        
        # Pattern 2: Question: ... Answer: ...
        if not training_data:
            pattern2 = r'Question[:\s]+(.+?)\s*Answer[:\s]+(.+?)(?=Question[:\s]+|$)'
            matches = re.findall(pattern2, text, re.DOTALL | re.IGNORECASE)
            for q, a in matches:
                q, a = q.strip(), a.strip()
                if q and a:
                    training_data.append({
                        "instruction": q,
                        "input": "",
                        "output": a
                    })
        
        # Pattern 3: Lines ending with ? followed by answer
        if not training_data:
            lines = text.split('\n')
            i = 0
            while i < len(lines):
                line = lines[i].strip()
                if line.endswith('?'):
                    question = line
                    answer_lines = []
                    i += 1
                    while i < len(lines) and not lines[i].strip().endswith('?'):
                        if lines[i].strip():
                            answer_lines.append(lines[i].strip())
                        i += 1
                    if answer_lines:
                        training_data.append({
                            "instruction": question,
                            "input": "",
                            "output": ' '.join(answer_lines)
                        })
                else:
                    i += 1
        
        return training_data
    
    def _extract_heading_content(self, content: List[Tuple]) -> List[Dict]:
        """Extract heading-content pairs."""
        training_data = []
        current_heading = None
        current_content = []
        
        for item_type, item_content in content:
            if item_type == 'heading':
                if current_heading and current_content:
                    training_data.append({
                        "instruction": f"Explain about: {current_heading}",
                        "input": "",
                        "output": '\n'.join(current_content)
                    })
                current_heading = item_content
                current_content = []
            elif item_type == 'text':
                current_content.append(item_content)
            elif item_type == 'table':
                # Format table as text
                table_text = self._format_table(item_content)
                current_content.append(table_text)
        
        # Don't forget last section
        if current_heading and current_content:
            training_data.append({
                "instruction": f"Explain about: {current_heading}",
                "input": "",
                "output": '\n'.join(current_content)
            })
        
        return training_data
    
    def _extract_chunks(self, content: List[Tuple], chunk_size: int = 500) -> List[Dict]:
        """Split content into chunks for summarization-style training."""
        training_data = []
        text = '\n'.join([c[1] if isinstance(c[1], str) else str(c[1]) for c in content])
        
        # Split into sentences
        sentences = re.split(r'(?<=[.!?])\s+', text)
        
        current_chunk = []
        current_length = 0
        
        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue
            
            if current_length + len(sentence) > chunk_size and current_chunk:
                # Save current chunk
                chunk_text = ' '.join(current_chunk)
                # Extract first few words as topic
                words = chunk_text.split()[:5]
                topic = ' '.join(words) + "..."
                
                training_data.append({
                    "instruction": f"Provide information about: {topic}",
                    "input": "",
                    "output": chunk_text
                })
                current_chunk = []
                current_length = 0
            
            current_chunk.append(sentence)
            current_length += len(sentence)
        
        # Don't forget last chunk
        if current_chunk:
            chunk_text = ' '.join(current_chunk)
            words = chunk_text.split()[:5]
            topic = ' '.join(words) + "..."
            training_data.append({
                "instruction": f"Provide information about: {topic}",
                "input": "",
                "output": chunk_text
            })
        
        return training_data
    
    # ==================== Helper Methods ====================
    
    def _is_heading(self, line: str) -> bool:
        """Check if a line is likely a heading."""
        # Markdown-style headings
        if line.startswith('#'):
            return True
        # Numbered headings
        if re.match(r'^\d+[\.\)]\s+\w', line):
            return True
        # ALL CAPS headings
        if line.isupper() and len(line) < 100:
            return True
        # Lines ending with colon
        if line.endswith(':') and len(line) < 80:
            return True
        # Underlined text (next line is ===)
        if re.match(r'^[=\-]{3,}$', line):
            return True
        return False
    
    def _clean_heading(self, line: str) -> str:
        """Clean heading text."""
        # Remove markdown symbols
        line = re.sub(r'^#+\s*', '', line)
        # Remove numbering
        line = re.sub(r'^\d+[\.\)]\s*', '', line)
        # Remove trailing colon
        line = line.rstrip(':')
        return line.strip()
    
    def _format_table(self, table_data: List[List[str]]) -> str:
        """Format table data as readable text."""
        if not table_data:
            return ""
        
        result = []
        headers = table_data[0] if table_data else []
        
        for row in table_data[1:]:
            row_text = []
            for i, cell in enumerate(row):
                if i < len(headers) and headers[i]:
                    row_text.append(f"{headers[i]}: {cell}")
                else:
                    row_text.append(cell)
            result.append(", ".join(row_text))
        
        return "\n".join(result)


# ==================== Colab Helper Functions ====================

def install_dependencies():
    """Install required dependencies in Colab."""
    import subprocess
    import sys
    
    packages = [
        'python-docx',
        'openpyxl',
        'pandas',
        'PyMuPDF',
        'python-pptx'
    ]
    
    for package in packages:
        subprocess.check_call([sys.executable, '-m', 'pip', 'install', '-q', package])
    
    print("✅ All dependencies installed!")


def upload_and_convert():
    """Upload files and convert to training data (for Colab)."""
    try:
        from google.colab import files
    except ImportError:
        print("This function is for Google Colab only.")
        return None
    
    print("📤 Upload your documents (DOCX, Excel, PDF, PPT, TXT)...")
    uploaded = files.upload()
    
    if not uploaded:
        print("No files uploaded.")
        return None
    
    converter = DocumentConverter()
    all_data = []
    
    for filename in uploaded.keys():
        try:
            data = converter.convert_file(filename)
            all_data.extend(data)
            print(f"✅ {filename}: {len(data)} training examples")
        except Exception as e:
            print(f"❌ {filename}: Error - {e}")
    
    print(f"\n📊 Total: {len(all_data)} training examples")
    return all_data


def combine_with_existing(new_data: List[Dict], existing_file: str) -> List[Dict]:
    """Combine new training data with existing JSON file."""
    try:
        with open(existing_file, 'r', encoding='utf-8') as f:
            existing_data = json.load(f)
        combined = existing_data + new_data
        print(f"✅ Combined: {len(existing_data)} existing + {len(new_data)} new = {len(combined)} total")
        return combined
    except FileNotFoundError:
        print(f"⚠️ {existing_file} not found. Using only new data.")
        return new_data


# ==================== CLI Interface ====================

if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(description="Convert documents to training data")
    parser.add_argument("files", nargs="+", help="Files to convert")
    parser.add_argument("-o", "--output", default="training_data.json", help="Output JSON file")
    parser.add_argument("-s", "--strategy", default="auto", 
                        choices=["auto", "faq", "heading", "chunk"],
                        help="Extraction strategy")
    
    args = parser.parse_args()
    
    converter = DocumentConverter()
    all_data = converter.convert_files(args.files, args.strategy)
    
    if all_data:
        converter.save_json(all_data, args.output)
    else:
        print("No training data extracted.")
